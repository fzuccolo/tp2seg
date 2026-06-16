from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.defensa import speaker_note_texts
from app.metricas import MetricResult


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
BORDER = RGBColor(219, 226, 239)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
WHITE = RGBColor(255, 255, 255)


def write_pptx(result: MetricResult, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, result)
    _slide_executive(prs, result)
    _slide_scope(prs, result)
    _slide_method(prs, result)
    _slide_chapters(prs, result)
    _slide_maturity_distribution(prs, result)
    _slide_top_gaps(prs, result)
    _slide_profile(prs, result)
    _slide_plan(prs, result)
    _slide_demo(prs, result)
    _slide_traceability(prs, result)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    _add_speaker_notes(output_path, speaker_note_texts(result))


def _blank_slide(prs: Presentation, title: str, kicker: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()

    if kicker:
        _textbox(slide, kicker.upper(), Inches(0.6), Inches(0.35), Inches(6.2), Inches(0.28), Pt(9), BLUE, bold=True)
    _textbox(slide, title, Inches(0.6), Inches(0.58), Inches(8.4), Inches(0.55), Pt(24), NAVY, bold=True)
    _textbox(
        slide,
        "TP2 Seguridad Informatica | Grupo 1 | TecnoHogar S.A. | ISO/IEC 27002:2022",
        Inches(0.6),
        Inches(7.05),
        Inches(8.5),
        Inches(0.25),
        Pt(8),
        MUTED,
    )
    return slide


def _slide_cover(prs: Presentation, result: MetricResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.45), 0, Inches(4.9), SLIDE_H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = BLUE
    panel.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.1), 0, Inches(0.18), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    _textbox(slide, "Tablero de Control\nde Seguridad", Inches(0.75), Inches(1.25), Inches(7.4), Inches(1.75), Pt(40), WHITE, bold=True)
    _textbox(slide, "TecnoHogar S.A. | ISO/IEC 27002:2022", Inches(0.8), Inches(3.2), Inches(6.6), Inches(0.36), Pt(17), RGBColor(219, 234, 254), bold=True)
    _textbox(slide, "Diagnostico, brechas y plan de mejora priorizado", Inches(0.8), Inches(3.85), Inches(6.4), Inches(0.34), Pt(15), RGBColor(226, 232, 240))
    _textbox(slide, "TP2 Seguridad Informatica", Inches(0.8), Inches(6.58), Inches(3.6), Inches(0.24), Pt(10), RGBColor(226, 232, 240))
    _textbox(slide, "Grupo 1", Inches(10.95), Inches(6.58), Inches(1.4), Inches(0.24), Pt(9), RGBColor(224, 242, 254), bold=True, align=PP_ALIGN.RIGHT)

    _metric_card(slide, "Controles", str(result.resumen["controles_evaluados"]), "universo evaluado", Inches(9.05), Inches(1.38), Inches(3.15), Inches(1.0), WHITE)
    _metric_card(slide, "Madurez", f"{result.resumen['madurez_global_pct']}%", "postura actual", Inches(9.05), Inches(2.72), Inches(3.15), Inches(1.0), WHITE)
    _metric_card(slide, "Plan", f"{result.resumen['proyectos']} proyectos", "acciones trazables", Inches(9.05), Inches(4.06), Inches(3.15), Inches(1.0), WHITE)


def _textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size: Pt,
    color: RGBColor = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    hyperlink: str | None = None,
    fit: bool = False,
):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    if hyperlink:
        run.hyperlink.address = hyperlink
        run.font.underline = True
    return shape


def _paragraph_box(slide, lines: list[str], x, y, w, h, title: str | None = None):
    box = _card(slide, x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.12)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Aptos"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
    for line in lines:
        p = tf.add_paragraph() if title or tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE
        p.space_after = Pt(3)
    return box


def _card(slide, x, y, w, h, fill: RGBColor = WHITE, line: RGBColor = BORDER, radius=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def _metric_card(slide, label: str, value: str, detail: str, x, y, w, h, color: RGBColor, value_link: str | None = None):
    card = _card(slide, x, y, w, h)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()

    compact = h < Inches(0.95)
    long_value = len(value) > 16
    value_size = Pt(18 if not compact else 14)
    if long_value:
        value_size = Pt(14 if not compact else 11)

    if compact:
        value_y, value_h = Inches(0.12), Inches(0.26)
        label_y, label_h = Inches(0.42), Inches(0.15)
        detail_y, detail_h = Inches(0.58), Inches(0.16)
    else:
        value_y, value_h = Inches(0.18), Inches(0.32)
        label_y, label_h = Inches(0.58), Inches(0.18)
        detail_y, detail_h = Inches(0.80), Inches(0.18)

    value_color = BLUE if value_link else NAVY
    _textbox(slide, value, x + Inches(0.18), y + value_y, w - Inches(0.32), value_h, value_size, value_color, bold=True, hyperlink=value_link, fit=True)
    _textbox(slide, label.upper(), x + Inches(0.18), y + label_y, w - Inches(0.32), label_h, Pt(7), SLATE, bold=True, fit=True)
    _textbox(slide, detail, x + Inches(0.18), y + detail_y, w - Inches(0.32), detail_h, Pt(7), MUTED, fit=True)
    return card


def _pill(slide, text: str, x, y, w, color: RGBColor):
    shape = _card(slide, x, y, w, Inches(0.32), fill=RGBColor(239, 246, 255), line=RGBColor(191, 219, 254))
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = color
    return shape


def _add_chart_title(slide, title: str, x, y, w):
    _textbox(slide, title, x, y, w, Inches(0.25), Pt(11), NAVY, bold=True)


def _style_chart(chart, max_scale: float | None = None, legend: bool = False):
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
    chart.has_title = False
    if hasattr(chart, "value_axis"):
        chart.value_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.color.rgb = SLATE
        chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE
        if max_scale is not None:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = max_scale
    if hasattr(chart, "category_axis"):
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.category_axis.tick_labels.font.color.rgb = SLATE
        chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
        chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE


def _series_color(chart, color: RGBColor, point_colors: list[RGBColor] | None = None):
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    if point_colors:
        for idx, point in enumerate(series.points):
            if idx >= len(point_colors):
                break
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = point_colors[idx]


def _bar_chart(
    slide,
    categories: list[str],
    values: list[float],
    x,
    y,
    w,
    h,
    title: str,
    color: RGBColor = BLUE,
    max_scale_floor: float | None = 100.0,
):
    _add_chart_title(slide, title, x, y - Inches(0.28), w)
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Valor", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, data).chart
    natural_max = max(values) * 1.12 if values else 1.0
    max_scale = max(max_scale_floor, natural_max) if max_scale_floor is not None else natural_max
    _style_chart(chart, max_scale=max_scale)
    _series_color(chart, color)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0.0'
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(8)
    plot.data_labels.font.color.rgb = SLATE
    return chart


def _column_chart(slide, categories: list[str], values: list[float], x, y, w, h, title: str, color: RGBColor = CYAN):
    _add_chart_title(slide, title, x, y - Inches(0.28), w)
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Controles", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data).chart
    _style_chart(chart, max_scale=max(values) * 1.2 if values else None)
    _series_color(chart, color)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(8)
    plot.data_labels.font.color.rgb = SLATE
    return chart


def _doughnut(slide, values: list[float], x, y, w, h, title: str):
    _add_chart_title(slide, title, x, y - Inches(0.28), w)
    data = CategoryChartData()
    data.categories = ["Madurez", "Brecha"]
    data.add_series("Postura", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, data).chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    _series_color(chart, BLUE, [GREEN, RED])
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0.0%'
    plot.data_labels.font.size = Pt(8)
    return chart


def _slide_executive(prs: Presentation, result: MetricResult) -> None:
    resumen = result.resumen
    slide = _blank_slide(prs, "Resumen ejecutivo", "Diagnostico")
    _textbox(
        slide,
        "93 controles ISO -> brechas ponderadas -> proyectos priorizados.\nCada accion queda trazable a evidencia y control.",
        Inches(0.6),
        Inches(1.18),
        Inches(7.1),
        Inches(0.58),
        Pt(11),
        SLATE,
    )
    _pill(slide, "ISO/IEC 27002:2022", Inches(9.6), Inches(0.7), Inches(1.6), BLUE)
    _pill(slide, "CMMI 0..5", Inches(11.35), Inches(0.7), Inches(1.1), BLUE)

    cards = [
        ("Madurez", f"{resumen['madurez_global_pct']}%", "promedio ponderado", BLUE),
        ("Brecha", f"{resumen['brecha_global_pct']}%", "distancia al objetivo", RED),
        ("Controles", str(resumen["controles_evaluados"]), "capitulos 5 a 8", CYAN),
        ("Proyectos", str(resumen["proyectos"]), "cartera priorizada", PURPLE),
        ("Quick wins", str(resumen["quick_wins"]), "impacto alto/esfuerzo bajo", GREEN),
        ("Jornadas", str(int(resumen["esfuerzo_total"])), "esfuerzo total", AMBER),
    ]
    x0, y0 = Inches(0.6), Inches(1.85)
    for idx, (label, value, detail, color) in enumerate(cards):
        x = x0 + Inches(2.2) * (idx % 3)
        y = y0 + Inches(1.42) * (idx // 3)
        _metric_card(slide, label, value, detail, x, y, Inches(2.0), Inches(1.18), color)

    _doughnut(
        slide,
        [float(resumen["madurez_global"]), float(resumen["brecha_global"])],
        Inches(7.65),
        Inches(1.65),
        Inches(4.95),
        Inches(3.2),
        "Madurez vs brecha global",
    )
    _paragraph_box(
        slide,
        [
            f"Capitulo mas debil: {resumen['capitulo_mas_debil']}",
            f"Capacidad mas debil: {resumen['capacidad_mas_debil']}",
            f"Proyecto prioritario: {resumen['proyecto_prioritario']}",
        ],
        Inches(7.75),
        Inches(5.08),
        Inches(4.7),
        Inches(1.25),
        "Lectura para direccion",
    )


def _slide_scope(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Alcance y modelo de datos", "Metodo")
    _paragraph_box(
        slide,
        [
            "TP1: inventario, clasificacion CID, procesos, aplicaciones y servidores.",
            "TP2: diagnostico ISO, brechas ponderadas, proyectos y trazabilidad.",
            "El mismo pipeline genera tablero, informe, slides HTML y PPTX.",
        ],
        Inches(0.6),
        Inches(1.35),
        Inches(4.25),
        Inches(2.0),
        "Continuidad TP1 -> TP2",
    )
    _paragraph_box(
        slide,
        [
            "Observado: TP1, roles y activos definidos en CSV.",
            "Definido: niveles CMMI, hallazgos, proyectos y vinculos.",
            "Derivado: madurez, brecha, quick wins y prioridades.",
        ],
        Inches(0.6),
        Inches(3.65),
        Inches(4.25),
        Inches(2.0),
        "Origen de los datos",
    )

    chapters = result.capitulos.copy().sort_values("capitulo")
    labels = [str(row["capitulo"]).split(" - ", 1)[0] for _, row in chapters.iterrows()]
    controls = [float(row["controles"]) for _, row in chapters.iterrows()]
    _column_chart(slide, labels, controls, Inches(5.35), Inches(1.55), Inches(6.9), Inches(3.0), "Controles evaluados por capitulo", BLUE)

    _textbox(
        slide,
        "Los cuatro capitulos evaluables de ISO/IEC 27002:2022 quedan cubiertos. La lectura puede ir desde un indicador ejecutivo hasta cada control individual.",
        Inches(5.55),
        Inches(4.9),
        Inches(6.3),
        Inches(0.72),
        Pt(12),
        SLATE,
    )
    _metric_card(slide, "Trazabilidad", "control -> evidencia -> proyecto", "cadena defendible ante preguntas", Inches(5.55), Inches(5.85), Inches(4.1), Inches(0.85), GREEN)


def _slide_chapters(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Madurez por capitulo ISO", "Resultados")
    chapters = result.capitulos.copy().sort_values("madurez_pct", ascending=True)
    labels = [str(value).replace("Controles ", "") for value in chapters["capitulo"]]
    values = [round(float(value), 1) for value in chapters["madurez_pct"]]
    _bar_chart(slide, labels, values, Inches(0.85), Inches(1.55), Inches(7.4), Inches(4.55), "Madurez porcentual por capitulo", BLUE)

    weakest = chapters.iloc[0]
    strongest = chapters.iloc[-1]
    _metric_card(slide, "Mas debil", f"{weakest['madurez_pct']:.1f}%", str(weakest["capitulo"]), Inches(8.75), Inches(1.55), Inches(3.65), Inches(1.05), RED)
    _metric_card(slide, "Mas fuerte", f"{strongest['madurez_pct']:.1f}%", str(strongest["capitulo"]), Inches(8.75), Inches(2.85), Inches(3.65), Inches(1.05), GREEN)
    _paragraph_box(
        slide,
        [
            "La brecha esta distribuida, pero tecnologia y personas quedan por debajo del promedio.",
            "Esto justifica un plan mixto: gobierno, hardening, vulnerabilidades, RRHH e incidentes.",
        ],
        Inches(8.75),
        Inches(4.25),
        Inches(3.65),
        Inches(1.4),
        "Interpretacion",
    )


def _slide_method(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Metodo de medicion", "Metodo")
    steps = [
        ("Control ISO", "catalogo comun"),
        ("Evidencia", "fuente del caso"),
        ("Nivel CMMI", "0 a 5"),
        ("Brecha", "1 - madurez"),
        ("Proyecto", "accion trazable"),
    ]
    x0, y = Inches(0.75), Inches(1.45)
    for idx, (title, subtitle) in enumerate(steps):
        x = x0 + Inches(2.45) * idx
        _card(slide, x, y, Inches(1.85), Inches(1.05), fill=WHITE, line=RGBColor(191, 219, 254))
        _textbox(slide, title, x + Inches(0.12), y + Inches(0.22), Inches(1.6), Inches(0.2), Pt(10), NAVY, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, subtitle, x + Inches(0.12), y + Inches(0.55), Inches(1.6), Inches(0.18), Pt(8), MUTED, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.83), y + Inches(0.38), Inches(0.52), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE
            arrow.line.fill.background()

    _paragraph_box(
        slide,
        [
            "Madurez: valor normalizado del nivel CMMI.",
            "Brecha: distancia entre situacion actual y objetivo.",
            "Brecha ponderada: brecha multiplicada por peso del control.",
            "Prioridad: brecha asociada + aporte esperado del proyecto.",
        ],
        Inches(0.85),
        Inches(3.35),
        Inches(5.4),
        Inches(2.25),
        "Definiciones de medicion",
    )
    _paragraph_box(
        slide,
        [
            "Observado: inventario, roles, activos y evidencia del caso.",
            "Definido: nivel CMMI, hallazgos, proyectos y vinculos.",
            "Derivado: madurez, brecha, quick wins y trazabilidad.",
        ],
        Inches(6.8),
        Inches(3.35),
        Inches(5.35),
        Inches(2.25),
        "Separacion de datos",
    )


def _slide_maturity_distribution(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Distribucion CMMI", "Resultados")
    dist = result.madurez_distribucion.copy()
    labels = [str(value).replace(" - ", "\n") for value in dist["madurez_nombre"]]
    values = [float(value) for value in dist["controles"]]
    _column_chart(slide, labels, values, Inches(0.8), Inches(1.55), Inches(7.15), Inches(4.6), "Cantidad de controles por nivel", CYAN)

    total_controls = int(sum(values)) or 1
    initial = int(dist.loc[dist["nivel_madurez"].isin([0, 1, 2]), "controles"].sum())
    defined = int(dist.loc[dist["nivel_madurez"].isin([3]), "controles"].sum())
    measured = int(dist.loc[dist["nivel_madurez"].isin([4, 5]), "controles"].sum())
    _metric_card(slide, "Inicial/Gestionado", f"{initial}", f"{initial / total_controls:.0%} del total", Inches(8.55), Inches(1.55), Inches(3.6), Inches(1.05), RED)
    _metric_card(slide, "Definido", f"{defined}", f"{defined / total_controls:.0%} del total", Inches(8.55), Inches(2.85), Inches(3.6), Inches(1.05), BLUE)
    _metric_card(slide, "Medido/Optimizado", f"{measured}", f"{measured / total_controls:.0%} del total", Inches(8.55), Inches(4.15), Inches(3.6), Inches(1.05), GREEN)
    _textbox(
        slide,
        "La lectura clave no es solamente el promedio: todavia hay muchos controles con practicas parciales y poca medicion formal.",
        Inches(8.65),
        Inches(5.55),
        Inches(3.35),
        Inches(0.8),
        Pt(11),
        SLATE,
    )


def _slide_profile(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Perfil de seguridad", "Lectura operacional")
    frame = result.capacidad_operacional.copy().head(8)
    labels = [str(value) for value in frame["atributo"]]
    values = [round(float(value), 1) for value in frame["madurez_pct"]]
    _bar_chart(slide, labels, values, Inches(0.85), Inches(1.55), Inches(7.4), Inches(4.75), "Madurez por capacidad operacional", PURPLE)

    weakest = frame.iloc[0]
    _metric_card(slide, "Capacidad mas debil", f"{float(weakest['madurez_pct']):.1f}%", str(weakest["atributo"]), Inches(8.75), Inches(1.55), Inches(3.65), Inches(1.05), RED)
    _paragraph_box(
        slide,
        [
            "Esta vista traduce controles ISO a capacidades de gestion.",
            "Para TecnoHogar, amenazas y vulnerabilidades quedan como alerta principal.",
            "La lectura apoya proyectos de hardening, monitoreo y SDLC seguro.",
        ],
        Inches(8.75),
        Inches(3.0),
        Inches(3.65),
        Inches(2.25),
        "Interpretacion",
    )


def _slide_top_gaps(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Brechas principales", "Prioridades")
    top = result.top_brechas.copy().head(8)
    labels = [f"{row.control_id}" for row in top.itertuples()]
    values = [round(float(value), 2) for value in top["peso_brecha"]]
    _bar_chart(
        slide,
        labels,
        values,
        Inches(0.75),
        Inches(1.55),
        Inches(6.6),
        Inches(4.7),
        "Top controles por brecha ponderada",
        RED,
        max_scale_floor=None,
    )

    table = top.loc[:, ["control_id", "control_nombre", "capitulo", "brecha_pct", "peso_brecha"]].head(5)
    x, y = Inches(7.65), Inches(1.35)
    _textbox(slide, "Lectura de las brechas", x, y, Inches(4.7), Inches(0.3), Pt(12), NAVY, bold=True)
    for idx, row in enumerate(table.itertuples(index=False), start=0):
        yy = y + Inches(0.52 + idx * 0.82)
        _card(slide, x, yy, Inches(4.65), Inches(0.62))
        _textbox(slide, str(row.control_id), x + Inches(0.12), yy + Inches(0.12), Inches(0.55), Inches(0.2), Pt(10), RED, bold=True)
        _textbox(slide, str(row.control_nombre), x + Inches(0.72), yy + Inches(0.09), Inches(2.95), Inches(0.22), Pt(8), NAVY, bold=True)
        detail = f"{row.capitulo} | brecha cruda {float(row.brecha_pct):.0f}%"
        _textbox(slide, detail, x + Inches(0.72), yy + Inches(0.34), Inches(2.95), Inches(0.18), Pt(7), MUTED)
        _textbox(slide, f"{float(row.peso_brecha):.2f}", x + Inches(3.82), yy + Inches(0.18), Inches(0.65), Inches(0.2), Pt(10), RED, bold=True, align=PP_ALIGN.RIGHT)


def _slide_plan(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Plan de accion priorizado", "Plan")
    projects = result.quick_wins.copy().head(10)

    _textbox(slide, "Matriz impacto / esfuerzo", Inches(0.75), Inches(1.12), Inches(4.8), Inches(0.28), Pt(12), NAVY, bold=True)
    matrix_x, matrix_y, matrix_w, matrix_h = Inches(0.75), Inches(1.76), Inches(5.35), Inches(4.38)
    _card(slide, matrix_x, matrix_y, matrix_w, matrix_h)
    mid_x = matrix_x + matrix_w / 2
    mid_y = matrix_y + matrix_h / 2
    for x in [mid_x]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, matrix_y, Inches(0.01), matrix_h)
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER
        line.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, matrix_x, mid_y, matrix_w, Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    _textbox(slide, "Quick wins", matrix_x + Inches(0.2), matrix_y + Inches(0.15), Inches(1.2), Inches(0.18), Pt(8), GREEN, bold=True)
    _textbox(slide, "Estrategicos", mid_x + Inches(0.2), matrix_y + Inches(0.15), Inches(1.3), Inches(0.18), Pt(8), BLUE, bold=True)
    _textbox(slide, "Esfuerzo ->", matrix_x + matrix_w - Inches(1.2), matrix_y + matrix_h + Inches(0.08), Inches(1.0), Inches(0.18), Pt(8), MUTED)
    _textbox(slide, "Prioridad ↑", matrix_x + Inches(0.08), matrix_y - Inches(0.26), Inches(1.05), Inches(0.18), Pt(8), MUTED)

    effort_max = float(projects["esfuerzo_jornadas"].max()) or 1.0
    priority_max = float(projects["prioridad"].max()) or 1.0
    color_by_quadrant = {"Quick win": GREEN, "Proyecto estrategico": BLUE, "Mejora tactica": AMBER, "Diferir": MUTED}
    for row in projects.itertuples(index=False):
        px = matrix_x + Inches(0.35) + (matrix_w - Inches(0.7)) * min(float(row.esfuerzo_jornadas) / effort_max, 1.0)
        py = matrix_y + matrix_h - Inches(0.35) - (matrix_h - Inches(0.7)) * min(float(row.prioridad) / priority_max, 1.0)
        color = color_by_quadrant.get(str(row.cuadrante), BLUE)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px - Inches(0.17), py - Inches(0.17), Inches(0.34), Inches(0.34))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = WHITE
        dot.line.width = Pt(1)
        _textbox(slide, str(row.proyecto_id), px + Inches(0.12), py - Inches(0.1), Inches(0.55), Inches(0.16), Pt(7), NAVY, bold=True)

    _textbox(slide, "Primeros frentes del plan", Inches(6.55), Inches(1.24), Inches(5.6), Inches(0.28), Pt(12), NAVY, bold=True)
    top = result.proyectos.copy().head(5)
    for idx, row in enumerate(top.itertuples(index=False), start=0):
        yy = Inches(1.62) + Inches(idx * 0.88)
        _card(slide, Inches(6.55), yy, Inches(5.75), Inches(0.68))
        _textbox(slide, str(row.proyecto_id), Inches(6.75), yy + Inches(0.15), Inches(0.7), Inches(0.2), Pt(10), BLUE, bold=True)
        _textbox(slide, str(row.titulo), Inches(7.45), yy + Inches(0.08), Inches(3.2), Inches(0.22), Pt(8), NAVY, bold=True)
        _textbox(slide, f"{row.plazo} | {int(row.esfuerzo_jornadas)} jornadas", Inches(7.45), yy + Inches(0.36), Inches(2.4), Inches(0.18), Pt(7), MUTED)
        _textbox(slide, f"{float(row.prioridad):.1f}", Inches(11.15), yy + Inches(0.19), Inches(0.75), Inches(0.2), Pt(11), GREEN, bold=True, align=PP_ALIGN.RIGHT)

    _textbox(
        slide,
        "Secuencia: quick wins de personas, incidentes y seguridad fisica; luego gobierno, IAM y SDLC; finalmente auditoria y mejora continua.",
        Inches(6.65),
        Inches(6.08),
        Inches(5.45),
        Inches(0.58),
        Pt(9),
        SLATE,
        fit=True,
    )


def _slide_demo(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Tablero publicado", "Entrega")
    _textbox(slide, "Vistas principales", Inches(0.75), Inches(1.28), Inches(4.4), Inches(0.3), Pt(13), NAVY, bold=True)
    steps = [
        ("Ejecutivo", "madurez, brecha y quick wins"),
        ("Mapa ISO", "93 controles y distribucion CMMI"),
        ("Brechas", "Pareto y concentracion de riesgo"),
        ("Plan", "proyectos, esfuerzo y roadmap"),
        ("Trazabilidad", "control -> proyecto -> evidencia"),
    ]
    for idx, (title, detail) in enumerate(steps):
        y = Inches(1.8 + idx * 0.8)
        _card(slide, Inches(0.8), y, Inches(5.1), Inches(0.56), fill=WHITE, line=BORDER)
        _textbox(slide, title, Inches(1.0), y + Inches(0.13), Inches(1.25), Inches(0.2), Pt(10), BLUE, bold=True)
        _textbox(slide, detail, Inches(2.25), y + Inches(0.13), Inches(3.35), Inches(0.2), Pt(9), SLATE)

    _paragraph_box(
        slide,
        [
            "Consulta ejecutiva del estado de seguridad.",
            "Navegacion desde indicadores hasta controles y proyectos.",
            "Base para seguimiento periodico del plan de mejora.",
        ],
        Inches(6.55),
        Inches(1.8),
        Inches(5.55),
        Inches(2.3),
        "Valor para gestion",
    )
    _metric_card(
        slide,
        "URL",
        "https://tp2seg.streamlit.app/",
        "tablero desplegado",
        Inches(6.55),
        Inches(4.55),
        Inches(4.8),
        Inches(0.9),
        CYAN,
        value_link="https://tp2seg.streamlit.app/",
    )


def _slide_traceability(prs: Presentation, result: MetricResult) -> None:
    slide = _blank_slide(prs, "Conclusiones y proximos pasos", "Conclusion")
    steps = [
        ("Control ISO", "93 metricas"),
        ("Evidencia", "TP1 + fuentes"),
        ("Brecha", "CMMI ponderado"),
        ("Proyecto", "plan priorizado"),
        ("Seguimiento", "tablero y CI"),
    ]
    x0, y = Inches(0.75), Inches(1.55)
    for idx, (title, subtitle) in enumerate(steps):
        x = x0 + Inches(2.45) * idx
        _card(slide, x, y, Inches(1.85), Inches(1.05), fill=WHITE, line=RGBColor(191, 219, 254))
        _textbox(slide, title, x + Inches(0.12), y + Inches(0.22), Inches(1.6), Inches(0.2), Pt(10), NAVY, bold=True, align=PP_ALIGN.CENTER)
        _textbox(slide, subtitle, x + Inches(0.12), y + Inches(0.55), Inches(1.6), Inches(0.18), Pt(8), MUTED, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.83), y + Inches(0.38), Inches(0.52), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE
            arrow.line.fill.background()

    _paragraph_box(
        slide,
        [
            "Cada indicador se vincula con control, evidencia, nivel CMMI y brecha.",
            "Cada proyecto se vincula con controles cubiertos, esfuerzo, prioridad y plazo.",
            "La generacion automatica permite actualizar tablero, informe y presentacion con cada cambio.",
        ],
        Inches(0.75),
        Inches(3.35),
        Inches(5.6),
        Inches(2.35),
        "Trazabilidad de decisiones",
    )
    _paragraph_box(
        slide,
        [
            "Usar el tablero como instrumento de gobierno, no solo como entrega academica.",
            "Aprobar quick wins y formalizar medicion periodica.",
            "Validar evidencias reales si TecnoHogar pasara de caso ficticio a auditoria.",
        ],
        Inches(6.75),
        Inches(3.35),
        Inches(5.4),
        Inches(2.35),
        "Plan inmediato",
    )


def _add_speaker_notes(pptx_path: Path, notes: list[str]) -> None:
    with zipfile.ZipFile(pptx_path, "r") as archive:
        files = {name: archive.read(name) for name in archive.namelist()}

    slide_count = len([name for name in files if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)])
    if slide_count != len(notes):
        raise ValueError(f"La cantidad de notas ({len(notes)}) no coincide con slides ({slide_count})")

    files["[Content_Types].xml"] = _with_content_type_overrides(files["[Content_Types].xml"].decode("utf-8"), slide_count).encode("utf-8")
    files["ppt/_rels/presentation.xml.rels"] = _with_notes_master_rel(files["ppt/_rels/presentation.xml.rels"].decode("utf-8")).encode("utf-8")
    files["ppt/notesMasters/notesMaster1.xml"] = _notes_master_xml().encode("utf-8")
    files["ppt/notesMasters/_rels/notesMaster1.xml.rels"] = _empty_relationships_xml().encode("utf-8")

    for index, note in enumerate(notes, start=1):
        slide_rels_path = f"ppt/slides/_rels/slide{index}.xml.rels"
        files[slide_rels_path] = _with_slide_notes_rel(files[slide_rels_path].decode("utf-8"), index).encode("utf-8")
        files[f"ppt/notesSlides/notesSlide{index}.xml"] = _notes_slide_xml(note).encode("utf-8")
        files[f"ppt/notesSlides/_rels/notesSlide{index}.xml.rels"] = _notes_slide_rels_xml(index).encode("utf-8")

    with zipfile.ZipFile(pptx_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, content in files.items():
            archive.writestr(name, content)


def _with_content_type_overrides(xml: str, slide_count: int) -> str:
    overrides = [
        '<Override PartName="/ppt/notesMasters/notesMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>'
    ]
    overrides.extend(
        f'<Override PartName="/ppt/notesSlides/notesSlide{index}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
        for index in range(1, slide_count + 1)
    )
    for override in overrides:
        part_name = re.search(r'PartName="([^"]+)"', override)
        if part_name and part_name.group(1) not in xml:
            xml = xml.replace("</Types>", f"{override}</Types>")
    return xml


def _with_notes_master_rel(xml: str) -> str:
    if "relationships/notesMaster" in xml:
        return xml
    rel_id = _next_rel_id(xml)
    rel = (
        f'<Relationship Id="{rel_id}" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" '
        'Target="notesMasters/notesMaster1.xml"/>'
    )
    return xml.replace("</Relationships>", f"{rel}</Relationships>")


def _with_slide_notes_rel(xml: str, index: int) -> str:
    if "relationships/notesSlide" in xml:
        return xml
    rel_id = _next_rel_id(xml)
    rel = (
        f'<Relationship Id="{rel_id}" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" '
        f'Target="../notesSlides/notesSlide{index}.xml"/>'
    )
    return xml.replace("</Relationships>", f"{rel}</Relationships>")


def _next_rel_id(xml: str) -> str:
    ids = [int(value) for value in re.findall(r'Id="rId(\d+)"', xml)]
    return f"rId{max(ids, default=0) + 1}"


def _empty_relationships_xml() -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
    )


def _notes_master_xml() -> str:
    return """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld name="">
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
    </p:spTree>
  </p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
</p:notesMaster>"""


def _notes_slide_xml(note: str) -> str:
    headings = {
        "Mensaje principal:",
        "Narrativa:",
        "Datos concretos:",
        "Puente:",
        "Preguntas esperables:",
        "Referencias:",
    }
    paragraphs = []
    for line in note.splitlines():
        text = line.strip()
        if not text:
            paragraphs.append("<a:p/>")
            continue
        paragraphs.append(_notes_paragraph_xml(text, text in headings))
    body = "".join(paragraphs)
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Speaker notes"/><p:cNvSpPr/><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr/>
        <p:txBody><a:bodyPr/><a:lstStyle/>{body}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>"""


def _notes_slide_rels_xml(index: int) -> str:
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{index}.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
</Relationships>"""


def _notes_paragraph_xml(text: str, heading: bool = False) -> str:
    if heading:
        return f'<a:p><a:r><a:rPr lang="es-AR" sz="1200" b="1"/><a:t>{_xml_escape(text)}</a:t></a:r></a:p>'

    parts = re.split(r"(\*\*[^*]+\*\*)", text)
    runs = []
    for part in parts:
        if not part:
            continue
        is_bold = part.startswith("**") and part.endswith("**")
        content = part[2:-2] if is_bold else part
        bold = ' b="1"' if is_bold else ""
        runs.append(f'<a:r><a:rPr lang="es-AR" sz="1200"{bold}/><a:t>{_xml_escape(content)}</a:t></a:r>')
    return f"<a:p>{''.join(runs)}</a:p>"


def _xml_escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
